from flask import Flask, Blueprint, render_template, request, redirect, url_for, session, jsonify, send_from_directory, flash
from werkzeug.utils import secure_filename
import os
import shutil
from pathlib import Path
from datetime import datetime
import google.generativeai as genai
import requests
from duckduckgo_search import DDGS
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
import PyPDF2
from pptx import Presentation
import re
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import time
import logging
import base64
from dotenv import load_dotenv
from passlib.hash import pbkdf2_sha256
import functools
import uuid
from supabase import create_client, Client

# Load environment variables
load_dotenv()

# Supabase setup - load from environment variables
SUPABASE_URL = os.getenv('SUPABASE_URL')
SUPABASE_KEY = os.getenv('SUPABASE_KEY')

if not SUPABASE_URL or not SUPABASE_KEY:
    print("Error: Supabase credentials not found in environment variables")
    print("Please set SUPABASE_URL and SUPABASE_KEY in your .env file")
    print("The application requires Supabase for user authentication")
    exit(1)  # Exit if no credentials are provided

# Initialize Supabase client
try:
    supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
    
    # Check if the 'users' table exists
    supabase.table('users').select('*').limit(1).execute()
    print("Successfully connected to Supabase")
except Exception as e:
    print(f"Error connecting to Supabase: {str(e)}")
    print("Make sure the 'users' table exists in your Supabase database with these columns:")
    print("user_id (text, primary key), name (text), email (text, unique), password (text), api_key (text), created_at (timestamp)")
    exit(1)  # Exit if Supabase connection fails

def find_user_by_email(email):
    """Find a user by email in Supabase"""
    try:
        response = supabase.table('users').select('*').eq('email', email).execute()
        
        # Check if we got a user
        if response.data and len(response.data) > 0:
            return response.data[0]
        return None
    except Exception as e:
        print(f"Error finding user in Supabase: {str(e)}")
        return None

def create_user(user_data):
    """Create a new user in Supabase"""
    try:
        response = supabase.table('users').insert(user_data).execute()
        
        # Check if the insert was successful
        if response.data and len(response.data) > 0:
            return True
        return False
    except Exception as e:
        print(f"Error creating user in Supabase: {str(e)}")
        return False

# Authentication middleware
def login_required(route_function):
    @functools.wraps(route_function)
    def wrapper(*args, **kwargs):
        if 'user_id' not in session:
            flash('Please log in to access this page.', 'error')
            return redirect(url_for('login'))
        return route_function(*args, **kwargs)
    return wrapper

# Initialize models
def initialize_genai_for_user(api_key=None):
    """Initialize Gemini API with the user's API key if provided, else with the default one"""
    if api_key:
        genai.configure(api_key=api_key)
    else:
        genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Constants and prompts
# Main prompt for summaries
PROMPT = """
You are a YouTube video summarizer designed for Mumbai University engineering students. 
Your task is to summarize the video transcript following the proper answer-writing format for 8-10 mark questions.

### **Instructions for Summarization:** 1. **Definition:** Start with a definition of the main topic and any closely related concepts.  
2. **Classification:** If the topic is broad, provide a **classification in a tree format** (use text-based representation like code blocks if needed).  
3. **Explanation:** Explain the topic in a structured, **stepwise or pointwise manner** to ensure clarity.  
4. **Diagrams:** If a diagram is necessary, Mention **"Draw a ____ Type of Diagram"** 
5. **Merits & Demerits:** List advantages and disadvantages **if applicable**.  
6. **Applications:** Mention real-world applications **if applicable**.    
7. **Conclusion:** End with a brief 2-3 line conclusion summarizing the key points.  
"""

# Concise prompt for shorter summaries
CONCISE_PROMPT = """
You are a YouTube video summarizer. Create a concise summary of the video in 5-10 key points.

Guidelines:
1. Each point should be clear and concise (1-2 lines max)
2. Use bullet points (•)
3. Focus on the most important concepts/ideas
4. Use keywords and technical terms where relevant
5. Keep the total summary within 200 words
6. Make points easy to remember and understand

Please provide a concise summary of this transcript:
"""

# Prompt for PDF and PPT documents
PDF_PPT_PROMPT = """
You are an educational content summarizer designed for engineering students. Analyze the provided content and create a comprehensive yet concise summary following this structure:

1. **Chapter Overview:**
   - Main topic and its significance
   - Key concepts covered
   - Prerequisites needed

2. **Topics Breakdown:**
   - List main topics and subtopics
   - Show relationships between concepts
   - Highlight important terms/definitions

3. **Simplified Explanations:**
   - Break down complex concepts
   - Use simple language
   - Provide examples where possible

4. **Key Points Summary:**
   - Bullet points of crucial information
   - Important formulas/equations (if any)
   - Common applications

5. **Study Focus:**
   - What to concentrate on
   - Potential exam topics
   - Common misconceptions to avoid

6. **Quick Revision Notes:**
   - 5-6 most important takeaways
   - Critical formulas/concepts to remember
   - Practice suggestion areas

Please analyze and summarize the following content:
"""

# Utility Functions
def format_content(text):
    """Convert markdown text to HTML for display in templates."""
    # Bold text
    text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'__(.*?)__', r'<strong>\1</strong>', text)
    
    # Italic text
    text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', text)
    text = re.sub(r'_(.*?)_', r'<em>\1</em>', text)
    
    # Headers
    text = re.sub(r'^# (.*?)$', r'<h1>\1</h1>', text, flags=re.MULTILINE)
    text = re.sub(r'^## (.*?)$', r'<h2>\1</h2>', text, flags=re.MULTILINE)
    text = re.sub(r'^### (.*?)$', r'<h3>\1</h3>', text, flags=re.MULTILINE)
    text = re.sub(r'^#### (.*?)$', r'<h4>\1</h4>', text, flags=re.MULTILINE)
    text = re.sub(r'^##### (.*?)$', r'<h5>\1</h5>', text, flags=re.MULTILINE)
    text = re.sub(r'^###### (.*?)$', r'<h6>\1</h6>', text, flags=re.MULTILINE)
    
    # Unordered lists
    # Replace bullet points with HTML list items
    lines = text.split('\n')
    in_list = False
    new_lines = []
    
    for line in lines:
        if line.strip().startswith('* ') or line.strip().startswith('- '):
            if not in_list:
                new_lines.append('<ul>')
                in_list = True
            new_lines.append(f'<li>{line.strip()[2:]}</li>')
        else:
            if in_list:
                new_lines.append('</ul>')
                in_list = False
            new_lines.append(line)
    
    if in_list:
        new_lines.append('</ul>')
    
    # Ordered lists
    text = '\n'.join(new_lines)
    lines = text.split('\n')
    in_list = False
    new_lines = []
    
    ordered_list_pattern = re.compile(r'^\d+\.\s')
    
    for line in lines:
        if ordered_list_pattern.match(line.strip()):
            if not in_list:
                new_lines.append('<ol>')
                in_list = True
            content = re.sub(r'^\d+\.\s', '', line.strip())
            new_lines.append(f'<li>{content}</li>')
        else:
            if in_list:
                new_lines.append('</ol>')
                in_list = False
            new_lines.append(line)
    
    if in_list:
        new_lines.append('</ol>')
    
    # Code blocks
    text = '\n'.join(new_lines)
    text = re.sub(r'```(.*?)```', r'<pre><code>\1</code></pre>', text, flags=re.DOTALL)
    
    # Inline code
    text = re.sub(r'`(.*?)`', r'<code>\1</code>', text)
    
    # Blockquotes
    text = re.sub(r'^> (.*?)$', r'<blockquote>\1</blockquote>', text, flags=re.MULTILINE)
    
    # Horizontal rule
    text = re.sub(r'^---$', r'<hr>', text, flags=re.MULTILINE)
    
    # Line breaks (ensure paragraphs are separated)
    paragraphs = text.split('\n\n')
    text = ''.join([f'<p>{p}</p>' if not (p.startswith('<') and p.endswith('>')) else p for p in paragraphs if p.strip()])
    
    return text

def extract_transcript(video_id):
    """Extract transcript from a YouTube video."""
    try:
        # First try with default 'en' language
        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            formatter = TextFormatter()
            return formatter.format_transcript(transcript)
        except Exception as primary_error:
            # If default language fails, try specific language variants
            try:
                # Try with 'en-IN' (English India) which is common
                transcript = YouTubeTranscriptApi.get_transcript(video_id, ['en-IN'])
                formatter = TextFormatter()
                return formatter.format_transcript(transcript)
            except:
                # If that fails too, try to get available languages and use the first one
                try:
                    transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                    # Get the first manual transcript or the first auto-generated one if no manual exists
                    for transcript in transcript_list:
                        # Try to get transcript and translate to English if needed
                        if transcript.is_translatable:
                            english_transcript = transcript.translate('en').fetch()
                            formatter = TextFormatter()
                            return formatter.format_transcript(english_transcript)
                        else:
                            # Use whatever transcript is available
                            transcript_data = transcript.fetch()
                            formatter = TextFormatter()
                            return formatter.format_transcript(transcript_data)
                except Exception as list_error:
                    # If all else fails, raise the original error
                    raise primary_error
                    
    except Exception as e:
        print(f"Error extracting transcript: {str(e)}")
        return f"Error: Could not extract transcript. {str(e)}"

def generate_summary(text, prompt):
    """Generate a summary using Google's Gemini API."""
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        combined_prompt = f"{prompt}\n\nText to summarize:\n{text}"
        
        response = model.generate_content(combined_prompt)
        
        if not response.text:
            return "Error: No response generated from AI model."
        
        return response.text
    except Exception as e:
        print(f"Error generating summary: {str(e)}")
        return f"Error generating summary: {str(e)}"

def determine_diagram_query(summary):
    """Determine a search query for relevant diagrams based on the summary."""
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        
        prompt = f"""
        Based on the following summary, provide a search query that would help find a relevant 
        educational diagram, chart, or image to complement the content. 
        The query should be concise (5-7 words) and focus on the main concept that would benefit 
        from visual explanation. Just provide the query text without any additional explanation.

        Summary:
        {summary}
        """
        
        response = model.generate_content(prompt)
        
        if not response.text:
            return "educational diagram"
        
        return response.text.strip()
    except Exception as e:
        print(f"Error determining diagram query: {str(e)}")
        return "educational diagram"

def download_relevant_images(query, num_images=3):
    """Download relevant images based on the query using DuckDuckGo search."""
    try:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(base_dir, "temp_images")
        os.makedirs(output_dir, exist_ok=True)
        
        images = []
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=num_images))
            
            for i, result in enumerate(results):
                try:
                    response = requests.get(result["image"], timeout=5)
                    if response.status_code == 200:
                        # Save the image
                        image_path = os.path.join(output_dir, f"image_{i}.jpg")
                        with open(image_path, "wb") as f:
                            f.write(response.content)
                        
                        # Store the relative path for use in templates
                        images.append(f"image_{i}.jpg")
                except Exception as e:
                    print(f"Error downloading image {i}: {str(e)}")
        
        return images
    except Exception as e:
        print(f"Error in download_relevant_images: {str(e)}")
        return f"Error downloading images: {str(e)}"

def extract_pdf_text(file_path):
    """Extract text from a PDF file."""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num in range(len(reader.pages)):
                text += reader.pages[page_num].extract_text() + "\n"
        return text
    except Exception as e:
        print(f"Error extracting PDF text: {str(e)}")
        return f"Error extracting PDF text: {str(e)}"

def extract_ppt_text(file_path):
    """Extract text from a PowerPoint presentation."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting PPT text: {str(e)}")
        return f"Error extracting PPT text: {str(e)}"

def take_full_page_screenshot(output_filename, html_content=None):
    """Take a screenshot of a full webpage using Selenium."""
    try:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(base_dir, "static", "screenshots")
        os.makedirs(output_dir, exist_ok=True)
        
        html_dir = os.path.join(base_dir, "temp_html")
        os.makedirs(html_dir, exist_ok=True)
        
        # Configure Chrome options
        chrome_options = Options()
        chrome_options.add_argument("--headless")
        chrome_options.add_argument("--window-size=1920,1080")
        chrome_options.add_argument("--disable-extensions")
        chrome_options.add_argument("--disable-gpu")
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--disable-dev-shm-usage")
        
        # Initialize driver
        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=chrome_options)
        
        try:
            if html_content:
                # Create a temporary HTML file
                temp_html_path = os.path.join(html_dir, f"temp_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
                with open(temp_html_path, 'w', encoding='utf-8') as f:
                    f.write(html_content)
                
                # Load local file
                driver.get(f"file:///{temp_html_path}")
            else:
                return {
                    'success': False,
                    'error': 'No HTML content provided for screenshot'
                }
            
            # Wait for page to load
            time.sleep(1)
            
            # Set viewport size to full content
            width = driver.execute_script("return Math.max(document.body.scrollWidth, document.documentElement.scrollWidth, document.documentElement.clientWidth);")
            height = driver.execute_script("return Math.max(document.body.scrollHeight, document.documentElement.scrollHeight, document.documentElement.clientHeight);")
            driver.set_window_size(width, height)
            
            # Take screenshot
            output_path = os.path.join(output_dir, output_filename)
            driver.save_screenshot(output_path)
            
            # Clean up temporary HTML file
            try:
                os.remove(temp_html_path)
            except:
                pass
            
            return {
                'success': True,
                'filename': output_filename
            }
            
        except Exception as e:
            print(f"Error taking screenshot: {str(e)}")
            
            # Try to take a screenshot of the error
            try:
                error_filename = f"error_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
                error_path = os.path.join(output_dir, error_filename)
                driver.save_screenshot(error_path)
                return {
                    'success': False,
                    'error': str(e),
                    'filename': error_filename
                }
            except:
                return {
                    'success': False,
                    'error': str(e)
                }
        finally:
            driver.quit()
            
    except Exception as e:
        print(f"Fatal error taking screenshot: {str(e)}")
        return {
            'success': False,
            'error': f"Fatal error: {str(e)}"
        }

# Initialize the application
app = Flask(__name__)
app.secret_key = os.urandom(24)
app.config['APPLICATION_NAME'] = 'InfoBite'

# Create title for HTML templates
@app.context_processor
def inject_app_name():
    return dict(app_name=app.config['APPLICATION_NAME'])

# Setup directories
base_dir = os.path.dirname(os.path.abspath(__file__))
temp_dir = Path(os.path.join(base_dir, "temp_images"))
try:
    shutil.rmtree(temp_dir)
except FileNotFoundError:
    pass
except Exception as e:
    print(f"Warning: Could not delete temp_images - {e}")

temp_dir.mkdir(exist_ok=True)

# Create the images directory if it doesn't exist
images_dir = os.path.join(base_dir, 'static', 'images')
if not os.path.exists(images_dir):
    os.makedirs(images_dir)

# Create screenshots directory if it doesn't exist
screenshots_dir = os.path.join(base_dir, 'static', 'screenshots')
if not os.path.exists(screenshots_dir):
    os.makedirs(screenshots_dir)

# Create uploads directory if it doesn't exist
uploads_dir = os.path.join(base_dir, 'uploads')
if not os.path.exists(uploads_dir):
    os.makedirs(uploads_dir)

# Register template filter
app.jinja_env.filters['format_content'] = format_content

# Main routes
@app.route('/', methods=['GET', 'POST'])
def home():
    """
    Home page route that handles form submissions.
    """
    # Check if user is authenticated
    if 'user_id' not in session:
        # If not authenticated, redirect to login page
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        youtube_link = request.form.get('youtube_link')
        prompt_option = request.form.get('prompt_option')
        custom_prompt = request.form.get('custom_prompt', '')

        if youtube_link:
            session['youtube_link'] = youtube_link
            session['prompt_option'] = prompt_option
            session['custom_prompt'] = custom_prompt
            return redirect(url_for('process_video'))
    return render_template('index.html')


@app.route('/about')
def about():
    """
    About page route.
    """
    # Check if user is authenticated - still allow access to about page
    # but with a noticeable notification to log in first
    return render_template('about.html')

# Video processing routes
@app.route('/process_video')
def process_video():
    """
    Process YouTube video and generate summary.
    """
    youtube_link = session.get('youtube_link')
    prompt_option = session.get('prompt_option')
    custom_prompt = session.get('custom_prompt', '')

    if not youtube_link:
        return "Error: No YouTube link provided."

    try:
        video_id = youtube_link.split("v=")[-1].split("&")[0]
        thumbnail_url = f"http://img.youtube.com/vi/{video_id}/0.jpg"

        transcript = extract_transcript(video_id)

        if "Error" in transcript:
            return render_template('result.html',
                                thumbnail_url=thumbnail_url,
                                summary=transcript,
                                images=[],
                                transcript="",
                                youtube_link=youtube_link)

        # Process with the appropriate prompt based on the option
        if prompt_option == "concise":
            summary = generate_summary(transcript, CONCISE_PROMPT)
            images = []
        elif prompt_option == "custom" and custom_prompt:
            summary = generate_summary(transcript, custom_prompt)
            images = []
        else:
            # Default to standard prompt
            summary = generate_summary(transcript, PROMPT)
            diagram_query = determine_diagram_query(summary)
            images = download_relevant_images(diagram_query)
            if isinstance(images, str):
                images = []

        if "Error" in summary:
            return render_template('result.html',
                                thumbnail_url=thumbnail_url,
                                summary=summary,
                                images=[],
                                transcript="",
                                youtube_link=youtube_link)

        # Store in session for later use
        session['summary'] = summary
        session['images'] = images
        session['current_image_index'] = 0
        
        return render_template('result.html',
                            thumbnail_url=thumbnail_url,
                            summary=summary,
                            images=images,
                            transcript=transcript,
                            youtube_link=youtube_link)
                            
    except Exception as e:
        error_msg = f"Error processing video: {str(e)}"
        print(error_msg)
        return render_template('result.html',
                            summary=error_msg,
                            images=[],
                            transcript="",
                            youtube_link=youtube_link)

@app.route('/static/temp_images/<path:filename>')
def serve_image(filename):
    """
    Serve image files from temp directory.
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    temp_images_dir = os.path.join(base_dir, 'temp_images')
    return send_from_directory(temp_images_dir, filename)

@app.route('/serve_image/<path:filename>')
def serve_image_direct(filename):
    """
    Serve image files from temp directory (alternate path).
    """
    base_dir = os.path.dirname(os.path.abspath(__file__))
    temp_images_dir = os.path.join(base_dir, 'temp_images')
    return send_from_directory(temp_images_dir, filename)

@app.route('/process_youtube', methods=['POST'])
def process_youtube():
    """
    Process YouTube link from form submission.
    """
    youtube_link = request.form.get('youtube_link')
    prompt_option = request.form.get('prompt_option')
    custom_prompt = request.form.get('custom_prompt', '')

    if youtube_link:
        session['youtube_link'] = youtube_link
        session['prompt_option'] = prompt_option
        session['custom_prompt'] = custom_prompt
        return redirect(url_for('process_video'))
    return redirect(url_for('home'))

@app.route('/ask_question', methods=['POST'])
def ask_question():
    """
    Answer questions based on document content.
    """
    try:
        data = request.json
        question = data.get('question')
        context = data.get('context')
        summary = data.get('summary')

        if not question or not context:
            return jsonify({
                'error': 'Missing question or context'
            }), 400

        prompt = f"""
        Based on the following document content and its summary, please answer the question.
        If the answer cannot be found in the content, say so.

        Document Summary:
        {summary}

        Full Document Content:
        {context}

        Question: {question}

        Please provide a clear and concise answer, using information from the document.
        If the information is not in the document, respond with "I cannot find information about this in the document."
        """

        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)

        if not response.text:
            return jsonify({
                'error': 'No response generated'
            }), 500

        return jsonify({
            'answer': response.text
        })

    except Exception as e:
        print(f"Error in ask_question: {str(e)}")
        return jsonify({
            'error': str(e)
        }), 500

@app.route('/video_screenshot', methods=['POST'])
def video_screenshot():
    """
    Take a screenshot of the video summary page using the current page HTML
    """
    try:
        # Get the HTML content from the request
        data = request.json
        html_content = data.get('html')
        
        if not html_content:
            return jsonify({
                'success': False,
                'error': 'No HTML content provided'
            })
        
        # Generate a filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"video_summary_{timestamp}.png"
        
        # Extract video ID from the HTML if possible
        try:
            video_id_match = html_content.split('id="video-id"')[1].split('>')[1].split('<')[0].strip()
            if video_id_match:
                filename = f"video_summary_{video_id_match}.png"
        except:
            # If extraction fails, just use the timestamp filename
            pass
        
        # Take the screenshot using the HTML content directly
        result = take_full_page_screenshot(output_filename=filename, html_content=html_content)
        
        # Check if screenshot was successful
        if result['success']:
            return jsonify({
                'success': True,
                'screenshot_path': f"/static/screenshots/{result['filename']}"
            })
        else:
            # Return the error screenshot if available
            if 'filename' in result:
                return jsonify({
                    'success': False,
                    'error': result['error'],
                    'screenshot_path': f"/static/screenshots/{result['filename']}"
                })
            else:
                return jsonify({
                    'success': False,
                    'error': result['error']
                })
    
    except Exception as e:
        print(f"Error taking screenshot: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        })

@app.route('/video_result')
def video_result():
    """
    Display video processing results (for screenshot purposes)
    """
    summary = session.get('summary')
    youtube_link = session.get('youtube_link')
    images = session.get('images', [])
    transcript = session.get('transcript', '')
    
    if not summary or not youtube_link:
        return redirect(url_for('home'))
    
    # Extract video ID for thumbnail
    video_id = youtube_link.split("v=")[-1].split("&")[0]
    thumbnail_url = f"http://img.youtube.com/vi/{video_id}/0.jpg"
    
    return render_template('result.html',
                           thumbnail_url=thumbnail_url,
                           summary=summary,
                           images=images,
                           transcript=transcript,
                           youtube_link=youtube_link)

# Document processing routes
@app.route('/process_document', methods=['POST'])
def process_document():
    """
    Process uploaded documents (PDF, PPT).
    """
    if 'document' not in request.files:
        return "No file uploaded"

    file = request.files['document']
    if file.filename == '':
        return "No file selected"

    try:
        # Get base directory
        base_dir = os.path.dirname(os.path.abspath(__file__))
        uploads_dir = os.path.join(base_dir, 'uploads')
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(uploads_dir, filename)

        os.makedirs(uploads_dir, exist_ok=True)

        file.save(file_path)

        if filename.endswith('.pdf'):
            text = extract_pdf_text(file_path)
        elif filename.endswith(('.ppt', '.pptx')):
            text = extract_ppt_text(file_path)
        else:
            return "Unsupported file format"

        if not text.strip():
            return "No text could be extracted from the document"

        summary = generate_summary(text, PDF_PPT_PROMPT)
        
        # Store the information in the session
        session['document_summary'] = summary
        session['document_text'] = text
        session['document_filename'] = filename

        os.remove(file_path)  # Clean up after processing

        return render_template('document_result.html',
                               summary=summary,
                               text=text)

    except Exception as e:
        print(f"Error processing document: {str(e)}")
        return f"Error processing document: {str(e)}"

@app.route('/document_screenshot', methods=['POST'])
def document_screenshot():
    """
    Take a screenshot of the document summary page using the current page HTML
    """
    try:
        # Get the HTML content from the request
        data = request.json
        html_content = data.get('html')
        
        if not html_content:
            return jsonify({
                'success': False,
                'error': 'No HTML content provided'
            })
        
        # Generate a filename with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"document_summary_{timestamp}.png"
        
        # Take the screenshot using the HTML content directly
        result = take_full_page_screenshot(output_filename=filename, html_content=html_content)
        
        # Check if screenshot was successful
        if result['success']:
            return jsonify({
                'success': True,
                'screenshot_path': f"/static/screenshots/{result['filename']}"
            })
        else:
            # Return the error screenshot if available
            if 'filename' in result:
                return jsonify({
                    'success': False,
                    'error': result['error'],
                    'screenshot_path': f"/static/screenshots/{result['filename']}"
                })
            else:
                return jsonify({
                    'success': False,
                    'error': result['error']
                })
    
    except Exception as e:
        print(f"Error taking screenshot: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        })

@app.route('/document_result')
def document_result():
    """
    Display document processing results (for screenshot purposes)
    """
    summary = session.get('document_summary')
    text = session.get('document_text')
    
    if not summary or not text:
        return redirect(url_for('home'))
    
    return render_template('document_result.html',
                           summary=summary,
                           text=text)

# Authentication Routes
@app.route('/login', methods=['GET', 'POST'])
def login():
    """Handle user login"""
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        
        # Validate inputs
        if not email or not password:
            flash('Please provide both email and password', 'error')
            return render_template('login.html')
        
        try:
            # Find user in database
            user = find_user_by_email(email)
            
            # Check if user exists and password is correct
            if user and pbkdf2_sha256.verify(password, user['password']):
                # Store user info in session
                session['user_id'] = user['user_id']
                session['user_name'] = user['name']
                session['user_email'] = user['email']
                session['user_api_key'] = user['api_key']
                
                # Initialize Gemini with user's API key
                initialize_genai_for_user(user['api_key'])
                
                flash('Login successful!', 'success')
                return redirect(url_for('home'))
            else:
                flash('Invalid email or password', 'error')
        except Exception as e:
            flash(f'Login error: {str(e)}', 'error')
    
    return render_template('login.html')

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    """Handle user registration"""
    if request.method == 'POST':
        name = request.form.get('name')
        email = request.form.get('email')
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')
        api_key = request.form.get('api_key')
        
        # Validate inputs
        if not name or not email or not password or not confirm_password or not api_key:
            flash('Please fill out all fields', 'error')
            return render_template('signup.html')
        
        if password != confirm_password:
            flash('Passwords do not match', 'error')
            return render_template('signup.html')
        
        try:
            # Check if email already exists
            existing_user = find_user_by_email(email)
            
            if existing_user:
                flash('Email already registered', 'error')
                return render_template('signup.html')
            
            # Validate API key
            try:
                test_genai = genai
                test_genai.configure(api_key=api_key)
                model = test_genai.GenerativeModel("gemini-1.5-flash")
                response = model.generate_content("Hello")
                if not response.text:
                    flash('API key validation failed', 'error')
                    return render_template('signup.html')
            except Exception as e:
                flash(f'Invalid API key: {str(e)}', 'error')
                return render_template('signup.html')
            
            # Create user data
            user_data = {
                'user_id': str(uuid.uuid4()),
                'name': name,
                'email': email,
                'password': pbkdf2_sha256.hash(password),
                'api_key': api_key,
                'created_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            }
            
            # Store in Supabase
            success = create_user(user_data)
            if success:
                flash('Account created successfully! Please log in.', 'success')
                return redirect(url_for('login'))
            else:
                flash('Error creating account. Please try again.', 'error')
                return render_template('signup.html')
            
        except Exception as e:
            flash(f'Error creating account: {str(e)}', 'error')
            return render_template('signup.html')
    
    return render_template('signup.html')

@app.route('/logout')
def logout():
    """Handle user logout"""
    # Clear session
    session.clear()
    flash('You have been logged out', 'success')
    return redirect(url_for('login'))

@app.route('/profile')
@login_required
def profile():
    """User profile page"""
    user_id = session.get('user_id')
    user_email = session.get('user_email')
    
    if not user_id or not user_email:
        session.clear()
        flash('User not found', 'error')
        return redirect(url_for('login'))
    
    # Get user from database
    user = find_user_by_email(user_email)
    
    if not user:
        session.clear()
        flash('User not found', 'error')
        return redirect(url_for('login'))
    
    return render_template('profile.html', user=user)

if __name__ == '__main__':
    app.run(debug=True)